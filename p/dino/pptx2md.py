from pptx import Presentation

def pptx_to_md(input_pptx, output_md):
    prs = Presentation(input_pptx)
    with open(output_md, 'w', encoding='utf-8') as md_file:
        for slide in prs.slides:
            md_file.write("# Slide\n\n")
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    md_file.write(f"{shape.text}\n\n")

pptx_to_md('DINO.pptx', 'DINO.md')
